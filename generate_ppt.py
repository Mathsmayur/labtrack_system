import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    # Create a new presentation object
    prs = Presentation()

    # Define some styles
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]

    # --- Slide 1: Title Slide ---
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "LabTrack System"
    subtitle.text = "Comprehensive Laboratory Management Solution\nProject Presentation"

    # --- Slide 2: Introduction ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Introduction & Overview"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "What is the LabTrack System?"
    
    p = tf.add_paragraph()
    p.text = "A full-stack web application designed to efficiently manage computer laboratories."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Empowers administrators to track PC hardware/software status in real-time."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Provides tools for scheduling, complaint resolution, and analytics."
    p.level = 1

    # --- Slide 3: Key Features ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Core Features"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Comprehensive feature set for complete lab control:"
    
    features = [
        "PC Management: Track status (Working, Defective) and location of all lab computers.",
        "Lab Scheduling: Manage timetables to prevent conflicts and allocate resources.",
        "Complaint System: Users can easily report hardware, software, or network issues.",
        "Notifications: Automated alerts keep users informed about their complaint status.",
        "Analytics Dashboard: Visual insights into lab health and frequent problem types."
    ]
    for feature in features:
        p = tf.add_paragraph()
        p.text = feature
        p.level = 1

    # --- Slide 4: Technology Stack ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Technology Stack"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Modern and robust technologies:"
    
    techs = [
        "Frontend: React.js (Vite) with custom modern UI/UX (Glassmorphism, dark themes).",
        "Backend: Java Spring Boot for a powerful, secure RESTful API.",
        "Database: MySQL 8.0 for reliable data storage.",
        "Security: Spring Security with JWT for authentication and role-based access.",
        "Containerization: Docker & Docker Compose for consistent deployment."
    ]
    for tech in techs:
        p = tf.add_paragraph()
        p.text = tech
        p.level = 1

    # --- Slide 5: User Roles & Workflows ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "User Roles & Workflows"
    body = slide.placeholders[1]
    tf = body.text_frame
    
    tf.text = "Admin Role"
    p = tf.add_paragraph()
    p.text = "Full access: manage labs, PCs, update schedules, view analytics, and resolve complaints."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "User Role (Student/Teacher)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "View lab schedules, report PC issues, and track the status of submitted complaints."
    p.level = 1

    # --- Slide 6: Deployment Strategy ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Deployment Strategy"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Ready for the Cloud:"
    
    deploy_points = [
        "Fully containerized architecture using Docker.",
        "docker-compose.yml orchestrates three microservices: frontend, backend, and database.",
        "Easy deployment to cloud providers like AWS.",
        "Environment variables used for secure and flexible configuration."
    ]
    for point in deploy_points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1

    # --- Slide 7: Conclusion & Future Scope ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Conclusion & Future Scope"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "LabTrack streamlines laboratory operations and maintenance."
    
    future = [
        "Future Enhancements:",
        "Integration with Active Directory/LDAP for centralized authentication.",
        "Mobile application for on-the-go access.",
        "Automated PC health checks via agent software."
    ]
    for item in future:
        p = tf.add_paragraph()
        p.text = item
        if "Enhancements" not in item:
            p.level = 1
        else:
            p.level = 0

    # --- Slide 8: Q&A ---
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Questions?"
    subtitle.text = "Thank you for your time."

    # Save the presentation
    output_path = "c:\\D\\project1\\LabTrack_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation successfully saved to: {output_path}")

if __name__ == "__main__":
    create_presentation()
